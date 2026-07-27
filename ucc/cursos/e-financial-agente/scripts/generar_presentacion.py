"""Genera presentacion PowerPoint profesional usando Claude para contenido expandido."""

import sys
import os
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import anthropic

def leer_brief(ruta):
    return Path(ruta).read_text(encoding="utf-8") if Path(ruta).exists() else None

def extraer_campo(texto, titulo):
    inicio = texto.find(f"## {titulo}")
    if inicio == -1:
        return ""
    inicio = texto.find("\n", inicio) + 1
    fin = texto.find("\n## ", inicio)
    if fin == -1:
        fin = len(texto)
    return texto[inicio:fin].strip()

def generar_contenido_claude(tema, objetivos, contenido, recursos):
    """Usa Claude para expandir y estructurar contenido profesionalmente"""
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        print("ERROR: ANTHROPIC_API_KEY no configurada")
        return None

    client = anthropic.Anthropic(api_key=api_key)

    prompt = f"""Genera un JSON valido con contenido educativo profesional.

TEMA: {tema}
OBJETIVOS: {objetivos}
CONTENIDO: {contenido}
RECURSOS: {recursos}

Responde SOLO con JSON valido (sin markdown). Usa este formato exactamente:
{{
  "resumen": "Texto de 1-2 lineas",
  "conceptos": ["concepto1", "concepto2", "concepto3", "concepto4", "concepto5"],
  "marco_teorico": "Parrafo sobre teoria. Otro parrafo. Otro parrafo.",
  "ejemplos": ["Ejemplo1 concreto", "Ejemplo2 del mundo real", "Ejemplo3 aplicable"],
  "actividades": ["Actividad1 practica", "Actividad2 interactiva", "Actividad3"],
  "lecturas": ["Libro o articulo", "Referencia", "URL o fuente"],
  "preguntas": ["Pregunta1?", "Pregunta2?", "Pregunta3?"]
}}

NO incluyas comillas dobles dentro de los valores. NO incluyas markdown."""

    message = client.messages.create(
        model="claude-opus-4-8",
        max_tokens=1500,
        messages=[{"role": "user", "content": prompt}]
    )

    try:
        texto = message.content[0].text.strip()
        print(f"DEBUG: Claude respondió con {len(texto)} caracteres")
        resultado = json.loads(texto)
        print(f"DEBUG: JSON parseado correctamente")
        return resultado
    except Exception as e:
        print(f"ERROR al parsear JSON: {e}")
        print(f"Texto recibido: {texto[:200]}")
        return None

def crear_presentacion(brief_path):
    """Crea presentacion con contenido generado por Claude"""
    brief = leer_brief(brief_path)
    if not brief:
        print(f"ERROR: No se pudo leer {brief_path}")
        return False

    tema = extraer_campo(brief, "Tema de la clase").split('\n')[0]
    semana_info = extraer_campo(brief, "Semana y fecha")
    objetivos = extraer_campo(brief, "Objetivos de aprendizaje")
    contenido = extraer_campo(brief, "Contenido clave a cubrir")
    recursos = extraer_campo(brief, "Recursos oficiales del temario")
    proximo_tema = extraer_campo(brief, "Proxima tema")

    fecha = ""
    for linea in semana_info.split('\n'):
        if "Fecha:" in linea:
            fecha = linea.split("Fecha:")[1].strip()
            break

    print(f"Generando contenido profesional con Claude...")
    contenido_expandido = generar_contenido_claude(tema, objetivos, contenido, recursos)

    if not contenido_expandido:
        print("ERROR: Claude no genero contenido")
        return False

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    COLOR_PRINCIPAL = RGBColor(0, 102, 204)
    COLOR_SECUNDARIO = RGBColor(51, 51, 51)
    COLOR_ACENTO = RGBColor(255, 102, 0)

    def agregar_titulo_slide(titulo, subtitulo=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_PRINCIPAL

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = titulo
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        if subtitulo:
            sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1.5))
            sf = sub.text_frame
            sf.word_wrap = True
            p = sf.paragraphs[0]
            p.text = subtitulo
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(200, 200, 200)

    def agregar_contenido_slide(titulo, items):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_PRINCIPAL
        bar.line.color.rgb = COLOR_PRINCIPAL

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = titulo
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(6))
        cf = content_box.text_frame
        cf.word_wrap = True

        for item in items:
            texto = str(item).strip().lstrip("-•*").strip()
            if texto:
                p = cf.add_paragraph()
                p.text = texto
                p.font.size = Pt(16)
                p.font.color.rgb = COLOR_SECUNDARIO
                p.space_before = Pt(8)
                p.space_after = Pt(8)

    # Slides
    agregar_titulo_slide(tema, f"Semana 1\n{fecha}")
    agregar_contenido_slide("Agenda", ["Objetivos de aprendizaje", "Contenido clave", "Actividades", "Recursos"])

    objetivos_list = [l.strip() for l in objetivos.split('\n') if l.strip()]
    agregar_contenido_slide("Objetivos", objetivos_list[:5])

    resumen = contenido_expandido.get("resumen", "")
    agregar_contenido_slide("Resumen Ejecutivo", [resumen])

    conceptos = contenido_expandido.get("conceptos", [])
    if isinstance(conceptos, list):
        agregar_contenido_slide("Conceptos Clave", conceptos[:6])

    marco = contenido_expandido.get("marco_teorico", [])
    if isinstance(marco, list):
        agregar_contenido_slide("Marco Teorico", marco[:4])
    elif isinstance(marco, str):
        agregar_contenido_slide("Marco Teorico", [marco])

    ejemplos = contenido_expandido.get("ejemplos", [])
    if isinstance(ejemplos, list):
        agregar_contenido_slide("Ejemplos Practicos", ejemplos[:4])

    actividades = contenido_expandido.get("actividades", [])
    if isinstance(actividades, list):
        agregar_contenido_slide("Actividades", actividades[:3])

    lecturas = contenido_expandido.get("lecturas", [])
    if isinstance(lecturas, list):
        agregar_contenido_slide("Lecturas Complementarias", lecturas[:5])

    preguntas = contenido_expandido.get("preguntas", [])
    if isinstance(preguntas, list):
        agregar_contenido_slide("Preguntas Clave", preguntas[:5])

    cierre = ["Revisamos los objetivos", "Discusion y preguntas", f"Proximo tema: {proximo_tema}"]
    agregar_contenido_slide("Cierre", cierre)

    output_dir = Path(brief_path).parent
    output_path = output_dir / "presentacion.pptx"
    prs.save(str(output_path))
    print(f"OK: presentacion.pptx generada")
    return True

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Uso: python generar_presentacion.py <ruta_brief>")
        sys.exit(1)
    crear_presentacion(sys.argv[1])

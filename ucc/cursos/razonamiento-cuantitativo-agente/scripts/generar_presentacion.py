"""
Genera presentación PowerPoint (.pptx) a partir del brief de clase.
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def leer_brief(ruta):
    """Lee el archivo brief.md"""
    if not Path(ruta).exists():
        return None
    return Path(ruta).read_text(encoding="utf-8")

def extraer_seccion(texto, titulo):
    """Extrae contenido de una sección del brief"""
    inicio = texto.find(f"## {titulo}")
    if inicio == -1:
        return ""
    inicio += len(f"## {titulo}") + 1
    fin = texto.find("\n## ", inicio)
    if fin == -1:
        fin = len(texto)
    return texto[inicio:fin].strip()

def crear_presentacion(brief_path):
    """Crea presentación PowerPoint desde el brief"""
    brief = leer_brief(brief_path)
    if not brief:
        print(f"ERROR: No se pudo leer {brief_path}")
        return False

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    COLOR_PRINCIPAL = RGBColor(0, 102, 204)
    COLOR_TEXTO = RGBColor(51, 51, 51)

    def agregar_slide_titulo(titulo, subtitulo=""):
        """Slide de título"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_PRINCIPAL

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = titulo
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        if subtitulo:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
            sub_frame = sub_box.text_frame
            p = sub_frame.paragraphs[0]
            p.text = subtitulo
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(220, 220, 220)

    def agregar_slide_contenido(titulo, contenido):
        """Slide de contenido con título y viñetas"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = titulo
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRINCIPAL

        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5.2))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for linea in contenido.split('\n'):
            linea = linea.strip()
            if linea.startswith('-'):
                linea = linea[1:].strip()
            if linea:
                p = content_frame.add_paragraph()
                p.text = linea
                p.font.size = Pt(18)
                p.font.color.rgb = COLOR_TEXTO

    tema = extraer_seccion(brief, "Tema de la clase").split('\n')[0] if "Tema de la clase" in brief else "Tema"
    objetivos = extraer_seccion(brief, "Objetivos de aprendizaje")
    contenido = extraer_seccion(brief, "Contenido clave a cubrir")
    semana = brief.split("Semana ")[1].split("\n")[0] if "Semana" in brief else "1"

    agregar_slide_titulo(tema, f"Semana {semana}")
    if objetivos:
        agregar_slide_contenido("Objetivos de Aprendizaje", objetivos)
    if contenido:
        agregar_slide_contenido("Contenido Clave", contenido)
    agregar_slide_contenido("Resumen", "- Conceptos clave cubiertos\n- Aplicaciones practicas\n- Proximo tema")

    output_dir = Path(brief_path).parent
    output_path = output_dir / "presentacion.pptx"
    prs.save(str(output_path))
    print(f"OK: presentacion.pptx generada")
    return True

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Uso: python generar_presentacion.py <ruta_brief>")
        sys.exit(1)
    crear_presentacion(sys.argv[1])

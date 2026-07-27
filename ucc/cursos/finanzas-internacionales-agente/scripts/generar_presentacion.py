"""Genera presentacion PowerPoint profesional usando Claude para contenido expandido."""

import sys
import os
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import anthropic

def leer_brief(ruta):
    return Path(ruta).read_text(encoding="utf-8") if Path(ruta).exists() else None

def extraer_campo(texto, titulo):
    inicio = texto.find(f"## {titulo}")
    if inicio == -1:
        return ""
    inicio = texto.find("\n", inicio) + 1
    fin = texto.find("\n## ", inicio)
    if fin == -1:
        fin = len(texto)
    return texto[inicio:fin].strip()

def generar_contenido_claude(tema, objetivos, contenido, recursos):
    """Usa Claude para expandir y estructurar contenido profesionalmente"""
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        print("ERROR: ANTHROPIC_API_KEY no configurada")
        return None

    client = anthropic.Anthropic(api_key=api_key)

    prompt = f"""DIRECTOR DE DISEÑO: Genera presentación con NARRATIVA FLUIDA, RIGOR ACADÉMICO y DISEÑO PROFESIONAL.

TEMA: {tema}
OBJETIVOS: {objetivos}
CONTENIDO: {contenido}
REFERENCIAS: {recursos}

REGLA DE ORO: Cada diapositiva abre con una FRASE ANCLA (afirmación completa ≤12 palabras). El contenido son 2-3 apoyos en ORACIONES LEGIBLES (no párrafos, no fragmentos). Las NOTAS DEL ORADOR llevan el guion hablado (4-8 líneas, conectores entre slides).

ESTRUCTURA JSON (RESPONDE SOLO JSON, SIN MARKDOWN):
{{
  "portada_frase_ancla": "Una afirmación sobre por qué importa el tema",
  "portada_notas": "4-8 líneas: gancho, impacto, transición a agenda",

  "marco_conceptual": {{
    "frase_ancla": "Definición breve y memorable (≤12 palabras)",
    "apoyos": ["Apoyo 1 (una oración legible)", "Apoyo 2", "Apoyo 3"],
    "notas_orador": "Explicación fluida con conectores, 4-8 líneas",
    "visual": "Descripción de ícono/diagrama/imagen a usar"
  }},

  "autores_clave": {{
    "frase_ancla": "Ej: 'Tres pensadores revolucionaron cómo entendemos esto'",
    "timeline": [
      {{"año": 1952, "autor": "Markowitz", "aporte": "Teoría Moderna de Portafolios", "impacto": "Base de diversificación"}},
      {{"año": 1995, "autor": "Scholes", "aporte": "Modelo de valuación", "impacto": "Herramienta profesional"}}
    ],
    "notas_orador": "Conexión histórica, evolución del pensamiento",
    "visual": "Línea de tiempo horizontal con nodos por autor"
  }},

  "caso_real_1": {{
    "frase_ancla": "Afirmación que resume la lección del caso",
    "contexto": "Dónde y cuándo (1 oración)",
    "hechos": "Qué pasó (1 oración con cifras si aplica)",
    "concepto": "Qué concepto del tema ilustra (1 oración)",
    "leccion": "Aprendizaje aplicable hoy (1 oración)",
    "notas_orador": "Narrativa de cómo ocurrió, por qué importa",
    "visual": "Foto/ícono/gráfico del caso"
  }},

  "caso_real_2": {{
    "frase_ancla": "Contraste o perspectiva diferente",
    "contexto": "...",
    "hechos": "...",
    "concepto": "...",
    "leccion": "...",
    "notas_orador": "...",
    "visual": "..."
  }},

  "datos_evidencia": {{
    "frase_ancla": "Ej: 'Los números confirman lo que la teoría predice'",
    "stats": [
      {{"cifra": "73%", "etiqueta": "Fondos que usan este modelo", "fuente": "Bloomberg 2024"}},
      {{"cifra": "$2.1T", "etiqueta": "AUM bajo estrategia relacionada", "fuente": "Vanguard Report 2023"}}
    ],
    "notas_orador": "Significado de los datos, contexto global",
    "visual": "Gráfico de barras o línea de tendencia"
  }},

  "aplicacion_ia": {{
    "frase_ancla": "IA y datos están transformando esta práctica",
    "herramienta_1": "Nombre/tecnología: uso en 1 oración",
    "herramienta_2": "...",
    "notas_orador": "Cómo cambia el trabajo hoy; ejemplos concretos",
    "visual": "Ícono de herramienta o diagrama de flujo"
  }},

  "debate_critico": {{
    "frase_ancla": "Aquí hay tensión: perspectivas en conflicto",
    "perspectiva_a": "Crítica o limitación (1-2 oraciones)",
    "perspectiva_b": "Defensa o contexto (1-2 oraciones)",
    "notas_orador": "Por qué ambas son válidas; cómo navegarlas",
    "visual": "Dos columnas en contraste o balanza"
  }},

  "actividad_clase": {{
    "frase_ancla": "Aplicar lo aprendido en 15 minutos",
    "pasos": ["Paso 1: ... (1 oración)", "Paso 2: ..."],
    "tiempo_badge": "15-20 min",
    "notas_orador": "Instrucciones detalladas y transición",
    "visual": "Ícono de actividad o cronómetro"
  }},

  "sintesis_5ideas": [
    {{"num": 1, "idea": "Idea clave 1 en 1 oración memorable (≤14 palabras)"}},
    {{"num": 2, "idea": "Idea clave 2"}},
    {{"num": 3, "idea": "Idea clave 3"}},
    {{"num": 4, "idea": "Idea clave 4"}},
    {{"num": 5, "idea": "Idea clave 5"}}
  ],

  "preguntas_profundas": [
    "¿Pregunta 1 de análisis crítico?",
    "¿Pregunta 2 de aplicación real?",
    "¿Pregunta 3 de impacto futuro?",
    "¿Pregunta 4 de perspectivas?",
    "¿Pregunta 5 de integración?"
  ],

  "referencias_apa": ["Referencia 1 APA 7", "Referencia 2", "Referencia 3 (mínimo 6)"],

  "narracion_general": "Resumen de cómo las diapositivas forman un arco: pregunta → desarrollo → casos → evidencia → síntesis → cierre."
}}

CALIDAD (CRÍTICO):
- Cada FRASE ANCLA es una afirmación completa, memorable y bien redactada (no un rótulo).
- Cada APOYO es una oración legible en español fluido (≤18 palabras), no fragmentos.
- NOTAS DEL ORADOR con conectores ("Ya vimos...", "Por eso...", "Esto significa...") para flujo narrativo.
- Datos reales con fuente y año. Casos específicos. Autores de verdad. SIN INVENTAR.
- Visuales descritos para cada sección (ícono, gráfico, línea de tiempo, tarjeta, etc.).
- El arco narrativo: cuando lees todas las frases ancla, cuentan una historia coherente."""

    message = client.messages.create(
        model="claude-opus-4-8",
        max_tokens=5500,
        messages=[{"role": "user", "content": prompt}]
    )

    try:
        texto = message.content[0].text.strip()
        print(f"DEBUG: Claude respondió con {len(texto)} caracteres")
        resultado = json.loads(texto)
        print(f"DEBUG: JSON parseado correctamente")
        return resultado
    except Exception as e:
        print(f"ERROR al parsear JSON: {e}")
        print(f"Texto recibido (primeros 500 chars):\n{texto[:500]}")
        print(f"Texto recibido (últimos 500 chars):\n{texto[-500:]}")
        return None

def crear_presentacion(brief_path):
    """Crea presentacion con contenido generado por Claude"""
    brief = leer_brief(brief_path)
    if not brief:
        print(f"ERROR: No se pudo leer {brief_path}")
        return False

    tema = extraer_campo(brief, "Tema de la clase").split('\n')[0]
    semana_info = extraer_campo(brief, "Semana y fecha")
    objetivos = extraer_campo(brief, "Objetivos de aprendizaje")
    contenido = extraer_campo(brief, "Contenido clave a cubrir")
    recursos = extraer_campo(brief, "Recursos oficiales del temario")
    proximo_tema = extraer_campo(brief, "Proxima tema")

    fecha = ""
    for linea in semana_info.split('\n'):
        if "Fecha:" in linea:
            fecha = linea.split("Fecha:")[1].strip()
            break

    print(f"Generando contenido profesional con Claude...")
    contenido_expandido = generar_contenido_claude(tema, objetivos, contenido, recursos)

    if not contenido_expandido:
        print("ERROR: Claude no genero contenido")
        return False

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    COLOR_PRINCIPAL = RGBColor(0, 102, 204)
    COLOR_SECUNDARIO = RGBColor(51, 51, 51)
    COLOR_ACENTO = RGBColor(255, 102, 0)

    def agregar_titulo_slide(titulo, subtitulo=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_PRINCIPAL

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = titulo
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        if subtitulo:
            sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1.5))
            sf = sub.text_frame
            sf.word_wrap = True
            p = sf.paragraphs[0]
            p.text = subtitulo
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(200, 200, 200)

    def agregar_contenido_slide(titulo, items, notas=""):
        """Crea slide con título, contenido en viñetas y notas del orador"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_PRINCIPAL
        bar.line.color.rgb = COLOR_PRINCIPAL

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = titulo
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(6))
        cf = content_box.text_frame
        cf.word_wrap = True

        for item in items:
            texto = str(item).strip().lstrip("-•*").strip()
            if texto:
                p = cf.add_paragraph()
                p.text = texto
                p.font.size = Pt(16)
                p.font.color.rgb = COLOR_SECUNDARIO
                p.space_before = Pt(8)
                p.space_after = Pt(8)

    def agregar_slide_frase_ancla(frase_ancla, apoyos, notas="", titulo_fallback=""):
        """Crea slide con frase ancla (bien redactada) y apoyos, con notas del orador"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_PRINCIPAL
        bar.line.color.rgb = COLOR_PRINCIPAL

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = frase_ancla if len(frase_ancla) <= 45 else titulo_fallback
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(6))
        cf = content_box.text_frame
        cf.word_wrap = True

        if isinstance(apoyos, list):
            for apoyo in apoyos:
                texto = str(apoyo).strip().lstrip("-•*").strip()
                if texto:
                    p = cf.add_paragraph()
                    p.text = texto
                    p.font.size = Pt(18)
                    p.font.color.rgb = COLOR_SECUNDARIO
                    p.space_before = Pt(8)
                    p.space_after = Pt(8)
        elif isinstance(apoyos, str):
            p = cf.add_paragraph()
            p.text = apoyos
            p.font.size = Pt(18)
            p.font.color.rgb = COLOR_SECUNDARIO

    # Portada
    agregar_titulo_slide(tema, f"Semana 1\n{fecha}")

    # Agenda
    agregar_contenido_slide("Agenda", ["Objetivos de aprendizaje", "Contenido clave", "Actividades", "Recursos"])

    # Objetivos
    objetivos_list = [l.strip() for l in objetivos.split('\n') if l.strip()]
    agregar_contenido_slide("Objetivos", objetivos_list[:5])

    # Marco conceptual (usa estructura de frase ancla)
    marco_block = contenido_expandido.get("marco_conceptual", {})
    if isinstance(marco_block, dict):
        frase = marco_block.get("frase_ancla", "Marco Conceptual")
        apoyos = marco_block.get("apoyos", [])
        agregar_slide_frase_ancla(frase, apoyos, "", "Marco Conceptual")
    elif isinstance(marco_block, str):
        agregar_contenido_slide("Marco Conceptual", [marco_block])

    # Autores clave (usa timeline si está disponible)
    autores_block = contenido_expandido.get("autores_clave", {})
    if isinstance(autores_block, dict):
        frase = autores_block.get("frase_ancla", "Autores Clave")
        timeline = autores_block.get("timeline", [])
        apoyos = [f"{a.get('autor', '')} ({a.get('año', '')}) — {a.get('aporte', '')}" for a in timeline if isinstance(timeline, list)]
        agregar_slide_frase_ancla(frase, apoyos, "", "Autores Clave")
    elif isinstance(autores_block, list):
        agregar_contenido_slide("Autores Clave", autores_block[:6])

    # Caso real 1 (usa estructura de tarjeta)
    caso1_block = contenido_expandido.get("caso_real_1", {})
    if isinstance(caso1_block, dict):
        frase = caso1_block.get("frase_ancla", "Caso Real 1")
        apoyos = [
            f"Contexto: {caso1_block.get('contexto', '')}",
            f"Hechos: {caso1_block.get('hechos', '')}",
            f"Concepto: {caso1_block.get('concepto', '')}",
            f"Lección: {caso1_block.get('leccion', '')}"
        ]
        agregar_slide_frase_ancla(frase, [a for a in apoyos if len(a) > 10], "", "Caso Real 1")
    elif isinstance(caso1_block, str):
        agregar_contenido_slide("Caso Real 1", [caso1_block])

    # Caso real 2
    caso2_block = contenido_expandido.get("caso_real_2", {})
    if isinstance(caso2_block, dict):
        frase = caso2_block.get("frase_ancla", "Caso Real 2")
        apoyos = [
            f"Contexto: {caso2_block.get('contexto', '')}",
            f"Hechos: {caso2_block.get('hechos', '')}",
            f"Concepto: {caso2_block.get('concepto', '')}",
            f"Lección: {caso2_block.get('leccion', '')}"
        ]
        agregar_slide_frase_ancla(frase, [a for a in apoyos if len(a) > 10], "", "Caso Real 2")

    # Datos y evidencia
    datos_block = contenido_expandido.get("datos_evidencia", {})
    if isinstance(datos_block, dict):
        frase = datos_block.get("frase_ancla", "Datos y Evidencia")
        stats = datos_block.get("stats", [])
        apoyos = [f"{s.get('cifra', '')} — {s.get('etiqueta', '')} ({s.get('fuente', '')})" for s in stats if isinstance(stats, list)]
        agregar_slide_frase_ancla(frase, apoyos, "", "Datos y Evidencia")
    elif isinstance(datos_block, str):
        agregar_contenido_slide("Datos y Evidencia", [datos_block])

    # Aplicación IA
    ia_block = contenido_expandido.get("aplicacion_ia", {})
    if isinstance(ia_block, dict):
        frase = ia_block.get("frase_ancla", "Aplicación de IA")
        herramientas = [
            ia_block.get("herramienta_1", ""),
            ia_block.get("herramienta_2", "")
        ]
        apoyos = [h for h in herramientas if h]
        agregar_slide_frase_ancla(frase, apoyos, "", "Aplicación de IA")
    elif isinstance(ia_block, str):
        agregar_contenido_slide("Aplicación de IA", [ia_block])

    # Debate crítico
    debate_block = contenido_expandido.get("debate_critico", {})
    if isinstance(debate_block, dict):
        frase = debate_block.get("frase_ancla", "Debate Crítico")
        apoyos = [
            f"Perspectiva A: {debate_block.get('perspectiva_a', '')}",
            f"Perspectiva B: {debate_block.get('perspectiva_b', '')}"
        ]
        agregar_slide_frase_ancla(frase, apoyos, "", "Debate Crítico")
    elif isinstance(debate_block, str):
        agregar_contenido_slide("Debate Crítico", [debate_block])

    # Actividad en clase
    actividad_block = contenido_expandido.get("actividad_clase", {})
    if isinstance(actividad_block, dict):
        frase = actividad_block.get("frase_ancla", "Actividad en Clase")
        pasos = actividad_block.get("pasos", [])
        agregar_slide_frase_ancla(frase, pasos, "", "Actividad en Clase")
    elif isinstance(actividad_block, str):
        agregar_contenido_slide("Actividad en Clase", [actividad_block])

    # Síntesis 5 ideas
    sintesis = contenido_expandido.get("sintesis_5ideas", [])
    if isinstance(sintesis, list):
        ideas = [f"{s.get('num', '')}. {s.get('idea', '')}" if isinstance(s, dict) else f"• {s}" for s in sintesis]
        agregar_contenido_slide("Síntesis: 5 Ideas Clave", ideas)

    # Preguntas profundas
    preguntas = contenido_expandido.get("preguntas_profundas", [])
    if isinstance(preguntas, list):
        agregar_contenido_slide("Preguntas Profundas", preguntas[:5])

    # Referencias APA
    referencias = contenido_expandido.get("referencias_apa", [])
    if isinstance(referencias, list):
        agregar_contenido_slide("Referencias APA 7", referencias[:8])

    # Cierre
    cierre = ["Revisamos los objetivos", "Discusión y preguntas", f"Próximo tema: {proximo_tema}"]
    agregar_contenido_slide("Cierre", cierre)

    output_dir = Path(brief_path).parent
    output_path = output_dir / "presentacion.pptx"
    prs.save(str(output_path))
    print(f"OK: presentacion.pptx generada")
    return True

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Uso: python generar_presentacion.py <ruta_brief>")
        sys.exit(1)
    crear_presentacion(sys.argv[1])

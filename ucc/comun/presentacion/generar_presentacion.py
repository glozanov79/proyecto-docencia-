#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
generar_presentacion.py  ·  Renderizador determinista de presentaciones UCC
---------------------------------------------------------------------------
El DISEÑO vive aquí, en código. El paso con LLM solo entrega DATOS (contenido.json);
este script decide colores, tipografía y layout. Así la identidad visual nunca
depende de que el modelo "recuerde" aplicarla: sale idéntica en todos los cursos.

Uso:
    python generar_presentacion.py contenido.json salida.pptx

contenido.json: ver contenido.schema.md para el contrato de cada tipo de diapositiva.
"""
import sys, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ══════════════ PALETA (reemplaza por los HEX OFICIALES de la UCC) ══════════════
PRIM   = RGBColor(0x0F, 0x6B, 0x4F)   # verde oscuro  · títulos, portada/cierre
ACC    = RGBColor(0x2F, 0xD6, 0xB0)   # aguamarina     · cifras, nodos, resaltados
APOYO  = RGBColor(0x7B, 0xD8, 0x8F)   # verde claro    · segundos elementos
TEXTO  = RGBColor(0x10, 0x22, 0x1B)   # tinta
SUTIL  = RGBColor(0x6B, 0x80, 0x78)   # gris neutro    · pies, fuentes
TINT   = RGBColor(0xEA, 0xF7, 0xF3)   # verde muy tenue· fondo de tarjetas
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
CLARO  = RGBColor(0xD8, 0xF3, 0xEA)   # texto claro sobre fondo oscuro
FONT   = "Arial"                       # cambia por la tipografía de la plantilla

ROUND = MSO_SHAPE.ROUNDED_RECTANGLE
RECT  = MSO_SHAPE.RECTANGLE
OVAL  = MSO_SHAPE.OVAL
W, H  = 13.333, 7.5

# ─────────────────────────────── helpers ───────────────────────────────
def _bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def _txt(slide, text, l, t, w, h, size=18, bold=False, color=TEXTO,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        fo = r.font; fo.size = Pt(size); fo.bold = bold; fo.name = font; fo.color.rgb = color
    return tb

def _shape(slide, kind, l, t, w, h, fill=None, line=None, line_w=1.0, radius=None):
    sp = slide.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    if radius is not None and kind == ROUND:
        try: sp.adjustments[0] = radius
        except Exception: pass
    sp.text_frame.paragraphs[0].text = ""
    return sp

def _label(sp, text, size, color, bold=True):
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.name = FONT; f.color.rgb = color

def _footer(slide, m):
    _txt(slide, f"{m['curso']} · {m['semana']}", 0.6, 7.05, 9, 0.3, size=10, color=SUTIL)
    _txt(slide, str(m["num"]), 12.2, 7.05, 0.6, 0.3, size=10, color=SUTIL, align=PP_ALIGN.RIGHT)

def _titulo(slide, titulo, frase=None):
    _txt(slide, titulo, 0.6, 0.45, 11.5, 0.8, size=32, bold=True, color=PRIM)
    if frase:
        _txt(slide, frase, 0.6, 1.35, 12.1, 0.7, size=23, color=TEXTO)

def _notas(slide, texto):
    if texto:
        slide.notes_slide.notes_text_frame.text = texto

# ─────────────────────────────── slide types ───────────────────────────────
def s_portada(slide, d, m):
    _bg(slide, PRIM)
    _shape(slide, OVAL, 10.7, -1.2, 3.6, 3.6, fill=RGBColor(0x15, 0x82, 0x63))
    _shape(slide, OVAL, 11.9, 5.6, 2.8, 2.8, fill=RGBColor(0x12, 0x78, 0x5B))
    _txt(slide, f"{m['curso']} · {m['semana']}".upper(), 0.75, 1.9, 10.5, 0.4,
         size=15, bold=True, color=ACC)
    _txt(slide, d.get("titulo", ""), 0.7, 2.5, 10.7, 2.0, size=44, bold=True, color=BLANCO)
    if d.get("subtitulo"):
        _txt(slide, d["subtitulo"], 0.75, 4.75, 11, 0.6, size=16, color=CLARO)
    _txt(slide, f"{m['institucion']} · {m['profesor']}", 0.75, 6.9, 12, 0.4,
         size=11, color=RGBColor(0xBF, 0xE6, 0xD8))

def s_lista_num(slide, d, m):
    _titulo(slide, d.get("titulo", "Agenda"), d.get("frase_ancla"))
    items = d.get("items") or d.get("objetivos") or []
    y = 2.7
    for i, it in enumerate(items[:6]):
        c = _shape(slide, OVAL, 0.7, y, 0.5, 0.5, fill=ACC); _label(c, str(i + 1), 16, PRIM)
        _txt(slide, it, 1.45, y, 11, 0.55, size=18, color=TEXTO, anchor=MSO_ANCHOR.MIDDLE)
        y += 0.78
    _footer(slide, m)

def s_tarjetas(slide, d, m):
    _titulo(slide, d.get("titulo", ""), d.get("frase_ancla"))
    cards = d.get("tarjetas", [])[:4]
    n = max(len(cards), 1); gap = 0.4; total = 12.1
    cw = (total - gap * (n - 1)) / n
    y, ch = 2.75, 3.0
    for i, card in enumerate(cards):
        x = 0.6 + i * (cw + gap)
        _shape(slide, ROUND, x, y, cw, ch, fill=TINT, radius=0.06)
        c = _shape(slide, OVAL, x + 0.32, y + 0.32, 0.7, 0.7, fill=ACC); _label(c, str(i + 1), 24, PRIM)
        _txt(slide, card.get("titulo", ""), x + 0.32, y + 1.2, cw - 0.64, 0.7,
             size=18, bold=True, color=PRIM)
        _txt(slide, card.get("texto", ""), x + 0.32, y + 1.9, cw - 0.64, ch - 2.0,
             size=13, color=TEXTO)
    _footer(slide, m)

def s_autores(slide, d, m):
    _titulo(slide, d.get("titulo", "Autores clave"), d.get("frase_ancla"))
    nodos = (d.get("nodos") or d.get("visual", {}).get("nodos") or [])[:6]
    n = max(len(nodos), 1)
    x0, x1, lineY = 2.0, 11.3, 4.55
    centros = [6.6] if n == 1 else [x0 + (x1 - x0) / (n - 1) * i for i in range(n)]
    if n > 1:
        _shape(slide, RECT, centros[0], lineY - 0.02, centros[-1] - centros[0], 0.04, fill=APOYO)
    for c, nd in zip(centros, nodos):
        _txt(slide, nd.get("anio", ""), c - 1.3, 3.75, 2.6, 0.5, size=22, bold=True,
             color=PRIM, align=PP_ALIGN.CENTER)
        _shape(slide, OVAL, c - 0.19, lineY - 0.19, 0.38, 0.38, fill=ACC, line=BLANCO, line_w=2)
        _txt(slide, nd.get("titulo", ""), c - 1.3, 4.95, 2.6, 0.4, size=15, bold=True,
             color=PRIM, align=PP_ALIGN.CENTER)
        _txt(slide, nd.get("aporte", ""), c - 1.35, 5.35, 2.7, 1.2, size=12,
             color=TEXTO, align=PP_ALIGN.CENTER)
    _footer(slide, m)

def _chart(slide, g, l, t, w, h):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
    cd = CategoryChartData()
    cd.categories = g.get("categorias", [])
    cd.add_series(g.get("serie", "Serie"), g.get("valores", []))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                Inches(l), Inches(t), Inches(w), Inches(h), cd)
    ch = gf.chart
    ch.has_legend = False
    ch.has_title = True
    ch.chart_title.text_frame.text = g.get("titulo", "")
    for p in ch.chart_title.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = PRIM; r.font.name = FONT
    plot = ch.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = g.get("num_format", '0"%"'); dl.number_format_is_linked = False
    dl.font.size = Pt(13); dl.font.color.rgb = PRIM; dl.font.name = FONT
    try: dl.position = XL_LABEL_POSITION.OUTSIDE_END
    except Exception: pass
    ser = plot.series[0]
    ser.format.fill.solid(); ser.format.fill.fore_color.rgb = ACC
    ca = ch.category_axis
    ca.tick_labels.font.size = Pt(12); ca.tick_labels.font.color.rgb = SUTIL
    ca.tick_labels.font.name = FONT
    try:
        ca.has_major_gridlines = False
        ch.value_axis.visible = False
        ch.value_axis.has_major_gridlines = False
    except Exception: pass

def s_datos(slide, d, m):
    _titulo(slide, d.get("titulo", "Datos y evidencia"), d.get("frase_ancla"))
    for i, st in enumerate(d.get("stats", [])[:3]):
        y = 2.4 + i * 1.5
        _txt(slide, st.get("cifra", ""), 0.6, y, 4.7, 0.7, size=36, bold=True, color=ACC)
        et = st.get("etiqueta", "")
        if st.get("fuente"): et += f" ({st['fuente']})"
        _txt(slide, et, 0.62, y + 0.72, 4.7, 0.7, size=12, color=SUTIL)
    g = d.get("grafico")
    if g:
        try: _chart(slide, g, 5.9, 2.4, 6.7, 4.2)
        except Exception: pass
    _footer(slide, m)

def s_caso(slide, d, m):
    _titulo(slide, d.get("titulo", "Caso real"), d.get("frase_ancla"))
    bloques = d.get("bloques", [])[:4]
    bx, by, bw, bh = [0.6, 6.85], [2.5, 4.75], 5.85, 2.0
    for i, b in enumerate(bloques):
        x, y = bx[i % 2], by[i // 2]
        _shape(slide, ROUND, x, y, bw, bh, fill=TINT, radius=0.05)
        c = _shape(slide, OVAL, x + 0.4, y + 0.55, 0.9, 0.9, fill=PRIM); _label(c, str(i + 1), 28, ACC)
        _txt(slide, b.get("etiqueta", "").upper(), x + 1.55, y + 0.35, bw - 1.8, 0.4,
             size=13, bold=True, color=ACC)
        _txt(slide, b.get("texto", ""), x + 1.55, y + 0.78, bw - 1.8, 1.1, size=13, color=TEXTO)
    _footer(slide, m)

def s_debate(slide, d, m):
    _titulo(slide, d.get("titulo", "Debate"), d.get("frase_ancla"))
    for col, x in [(d.get("columna_a", {}), 0.6), (d.get("columna_b", {}), 6.85)]:
        w, y, hh = 5.85, 2.7, 3.9
        _shape(slide, ROUND, x, y, w, hh, fill=TINT, radius=0.05)
        _txt(slide, col.get("titulo", ""), x + 0.4, y + 0.35, w - 0.8, 0.6,
             size=18, bold=True, color=PRIM)
        pts = ["•  " + p for p in col.get("puntos", [])]
        _txt(slide, pts, x + 0.4, y + 1.1, w - 0.8, hh - 1.3, size=14, color=TEXTO)
    _footer(slide, m)

def s_actividad(slide, d, m):
    _titulo(slide, d.get("titulo", "Actividad en clase"), d.get("frase_ancla"))
    if d.get("tiempo"):
        b = _shape(slide, ROUND, 10.5, 0.5, 2.2, 0.65, fill=ACC, radius=0.5)
        _label(b, d["tiempo"], 16, PRIM)
    y = 2.6
    for i, p in enumerate(d.get("pasos", [])[:6]):
        c = _shape(slide, OVAL, 0.7, y, 0.55, 0.55, fill=PRIM); _label(c, str(i + 1), 18, ACC)
        _txt(slide, p, 1.5, y, 10.9, 0.6, size=15, color=TEXTO, anchor=MSO_ANCHOR.MIDDLE)
        y += 0.8
    _footer(slide, m)

def s_sintesis(slide, d, m):
    _titulo(slide, d.get("titulo", "Síntesis"), d.get("frase_ancla"))
    ideas = d.get("ideas", [])[:5]
    y, rh, gap = 2.5, 0.72, 0.12
    for i, idea in enumerate(ideas):
        _shape(slide, ROUND, 0.6, y, 12.1, rh, fill=TINT, radius=0.25)
        c = _shape(slide, OVAL, 0.8, y + 0.11, 0.5, 0.5, fill=ACC); _label(c, str(i + 1), 18, PRIM)
        _txt(slide, idea, 1.6, y, 10.9, rh, size=15, color=TEXTO, anchor=MSO_ANCHOR.MIDDLE)
        y += rh + gap
    _footer(slide, m)

def s_cierre(slide, d, m):
    _bg(slide, PRIM)
    _txt(slide, d.get("titulo", "Cierre"), 0.7, 2.3, 11, 1.2, size=40, bold=True, color=BLANCO)
    if d.get("frase_ancla"):
        _txt(slide, d["frase_ancla"], 0.75, 3.65, 11, 0.8, size=22, color=CLARO)
    if d.get("proximo_tema"):
        _txt(slide, "PRÓXIMO TEMA", 0.75, 4.9, 11, 0.4, size=13, bold=True, color=ACC)
        _txt(slide, d["proximo_tema"], 0.75, 5.3, 11, 0.9, size=20, color=BLANCO)

def s_referencias(slide, d, m):
    _titulo(slide, d.get("titulo", "Referencias"))
    items = d.get("items", [])
    half = (len(items) + 1) // 2 if len(items) > 5 else len(items)
    c1 = [f"{i+1}. {r}" for i, r in enumerate(items[:half])]
    c2 = [f"{half+i+1}. {r}" for i, r in enumerate(items[half:])]
    _txt(slide, c1, 0.6, 2.0, 6.0, 4.8, size=11, color=TEXTO)
    if c2:
        _txt(slide, c2, 6.9, 2.0, 6.0, 4.8, size=11, color=TEXTO)
    _footer(slide, m)

def s_generico(slide, d, m):
    _titulo(slide, d.get("titulo", d.get("tipo", "")), d.get("frase_ancla"))
    apoyos = d.get("apoyos") or d.get("items") or []
    if apoyos:
        _txt(slide, ["•  " + a for a in apoyos], 0.6, 2.6, 12, 4, size=18, color=TEXTO)
    _footer(slide, m)

RENDER = {
    "portada": s_portada, "agenda": s_lista_num, "objetivos": s_lista_num,
    "autores": s_autores, "timeline": s_autores, "marco": s_tarjetas,
    "aplicacion_ia": s_tarjetas, "datos": s_datos, "caso": s_caso,
    "debate": s_debate, "actividad": s_actividad, "sintesis": s_sintesis,
    "cierre": s_cierre, "referencias": s_referencias,
}

def main():
    if len(sys.argv) < 2:
        print("Uso: python generar_presentacion.py contenido.json [salida.pptx]"); sys.exit(1)
    inp = sys.argv[1]
    out = sys.argv[2] if len(sys.argv) > 2 else "salida.pptx"
    data = json.load(open(inp, encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = Inches(W); prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    meta_base = {
        "curso": data.get("curso", ""), "institucion": data.get("institucion", "Universidad Cooperativa de Colombia"),
        "profesor": data.get("profesor", ""),
    }
    for i, d in enumerate(data.get("diapositivas", []), start=1):
        slide = prs.slides.add_slide(blank)
        m = dict(meta_base); m["semana"] = data.get("semana", ""); m["num"] = i
        fn = RENDER.get(d.get("tipo"), s_generico)
        try:
            fn(slide, d, m)
        except Exception as e:
            _titulo(slide, d.get("titulo", d.get("tipo", "diapositiva")))
            _txt(slide, f"[error de render: {e}]", 0.6, 2.6, 12, 1, size=12, color=SUTIL)
        _notas(slide, d.get("notas_orador"))
    prs.save(out)
    print(f"OK · {len(data.get('diapositivas', []))} diapositivas → {out}")

if __name__ == "__main__":
    main()
